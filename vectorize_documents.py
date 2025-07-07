import os
from azure.storage.blob import BlobServiceClient
from openai import AzureOpenAI
from dotenv import load_dotenv
import numpy as np
from tqdm import tqdm
import json
from datetime import datetime
import PyPDF2
import io
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

class DocumentVectorizer:
    def __init__(self):
        # Load environment variables
        load_dotenv()
        
        # Initialize Azure OpenAI client
        self.openai_client = AzureOpenAI(
            api_key=os.getenv('AZURE_API_KEY'),
            api_version="2024-02-15-preview",
            azure_endpoint=os.getenv('AZURE_ENDPOINT')
        )
        
        # Initialize Azure Blob Storage client
        self.blob_service_client = BlobServiceClient.from_connection_string(
            os.getenv('AZURE_CONNECTION_STRING')
        )
        
        # Embedding model name
        self.embedding_model = "text-embedding-ada-002"
        
        # Create vectors directory if it doesn't exist
        self.vectors_dir = "vectors"
        os.makedirs(self.vectors_dir, exist_ok=True)

    def get_embedding(self, text):
        """Get embedding for a single text using Azure OpenAI."""
        try:
            response = self.openai_client.embeddings.create(
                input=text,
                model=self.embedding_model
            )
            return response.data[0].embedding
        except Exception as e:
            print(f"Error getting embedding: {str(e)}")
            return None

    def extract_text_from_pdf(self, pdf_bytes):
        """Extract text from PDF bytes."""
        try:
            pdf_file = io.BytesIO(pdf_bytes)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            
            # Extract text from each page
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from PDF: {str(e)}")
            return None

    def extract_text_from_pptx(self, pptx_bytes):
        """Extract text from PPTX bytes."""
        try:
            pptx_file = io.BytesIO(pptx_bytes)
            prs = Presentation(pptx_file)
            text = ""
            
            # Extract text from each slide
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from PPTX: {str(e)}")
            return None

    def extract_text_from_docx(self, docx_bytes):
        try:
            docx_file = io.BytesIO(docx_bytes)
            doc = Document(docx_file)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from DOCX: {str(e)}")
            return None

    def extract_text_from_image(self, image_bytes):
        try:
            image = Image.open(io.BytesIO(image_bytes))
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            print(f"Error extracting text from image: {str(e)}")
            return None

    def read_blob_content(self, blob_client):
        """Read content from a blob."""
        try:
            dow_stream = blob_client.download_blob()
            content = dow_stream.readall()
            name = blob_client.blob_name.lower()
            if name.endswith('.pdf'):
                return self.extract_text_from_pdf(content)
            elif name.endswith('.pptx'):
                return self.extract_text_from_pptx(content)
            elif name.endswith('.docx'):
                return self.extract_text_from_docx(content)
            elif name.endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.webp')):
                return self.extract_text_from_image(content)
            else:
                try:
                    return content.decode('utf-8')
                except UnicodeDecodeError:
                    print(f"Warning: Could not decode {blob_client.blob_name} as text. Skipping.")
                    return None
        except Exception as e:
            print(f"Error reading blob {blob_client.blob_name}: {str(e)}")
            return None

    def split_text(self, text, max_chunk_size=2000):
        """Split text into chunks of max_chunk_size characters."""
        return [text[i:i+max_chunk_size] for i in range(0, len(text), max_chunk_size)]

    def process_container(self, container_name):
        """Process all documents in a container."""
        print(f"\nProcessing container: {container_name}")
        
        # Get container client
        container_client = self.blob_service_client.get_container_client(container_name)
        
        # List all blobs in the container
        blob_list = list(container_client.list_blobs())
        
        # Create a list to store document data
        documents = []
        
        # Process each blob
        for blob in tqdm(blob_list, desc="Processing documents"):
            # Skip non-PDF and non-PPTX files
            if not (blob.name.lower().endswith('.pdf') or blob.name.lower().endswith('.pptx')):
                continue
                
            blob_client = container_client.get_blob_client(blob.name)
            
            # Read blob content
            content = self.read_blob_content(blob_client)
            if not content:
                continue
            
            # Split content into chunks
            chunks = self.split_text(content, max_chunk_size=2000)
            for chunk_idx, chunk in enumerate(chunks):
                # Get embedding
                embedding = self.get_embedding(chunk)
                if not embedding:
                    continue
                # Store document data for each chunk
                documents.append({
                    'blob_name': blob.name,
                    'container': container_name,
                    'content': chunk,
                    'embedding': embedding,
                    'last_modified': blob.last_modified.isoformat(),
                    'size': blob.size,
                    'chunk_index': chunk_idx,
                    'num_chunks': len(chunks)
                })
        
        return documents

    def save_vectors(self, documents, container_name):
        """Save vectors and metadata to files."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Save embeddings as numpy array
        embeddings = np.array([doc['embedding'] for doc in documents])
        np.save(f"{self.vectors_dir}/{container_name}_embeddings_{timestamp}.npy", embeddings)
        
        # Save metadata as JSON with content
        metadata = [{
            'blob_name': doc['blob_name'],
            'container': doc['container'],
            'last_modified': doc['last_modified'],
            'size': doc['size'],
            'content': doc['content']  # Store the content for snippet generation
        } for doc in documents]
        
        with open(f"{self.vectors_dir}/{container_name}_metadata_{timestamp}.json", 'w') as f:
            json.dump(metadata, f, indent=2)

    def vectorize_all_containers(self):
        """Process all containers and create vector embeddings."""
        try:
            # List all containers
            containers = self.blob_service_client.list_containers()
            
            for container in containers:
                # Process container
                documents = self.process_container(container.name)
                
                if documents:
                    # Save vectors and metadata
                    self.save_vectors(documents, container.name)
                    print(f"Successfully processed {len(documents)} documents in container {container.name}")
                else:
                    print(f"No documents were processed in container {container.name}")
                
        except Exception as e:
            print(f"Error during vectorization: {str(e)}")

def main():
    vectorizer = DocumentVectorizer()
    vectorizer.vectorize_all_containers()

if __name__ == "__main__":
    main() 